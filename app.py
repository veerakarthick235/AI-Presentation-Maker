import os
import json
import google.generativeai as genai
from flask import Flask, request, jsonify, render_template, url_for
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from gtts import gTTS
import requests
from io import BytesIO
from bs4 import BeautifulSoup

# --- Configuration & Setup ---

load_dotenv()
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
genai.configure(api_key=GOOGLE_API_KEY)
app = Flask(__name__)
OUTPUT_FOLDER = os.path.join('static', 'outputs')
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


# --- Helper Functions (Shared) ---

def clean_json_response(text):
    """
    Extracts JSON from Gemini response safely.
    """
    try:
        start_index = text.find('```json')
        end_index = text.rfind('```')

        if start_index != -1 and end_index != -1:
            json_str = text[start_index + 7:end_index].strip()
        else:
            json_str = text.strip()

        return json.loads(json_str)
    except Exception:
        print("⚠️ JSON decoding failed. Raw response:")
        print(text)
        return None


def generate_speech(text, speech_filename):
    """
    Converts text to speech and saves MP3.
    """
    print(f"Generating speech for: {speech_filename}")
    try:
        tts = gTTS(text=text, lang='en')
        speech_path = os.path.join(app.config['OUTPUT_FOLDER'], speech_filename)
        tts.save(speech_path)
        return speech_path
    except Exception as e:
        print(f"Error generating speech: {e}")
        return None


def create_powerpoint(slides_data):
    """
    Creates PowerPoint with AI-generated text and audio.
    """
    print("🧩 Creating PowerPoint...")
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = slides_data[0]['title']
    slide.placeholders[1].text = "AI-Generated Presentation"

    # Content slides
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data['title']

        speech_filename = f"slide_{i + 1}.mp3"
        speech_path = generate_speech(slide_data['content'], speech_filename)

        slide.placeholders[1].text = slide_data['content']
        if speech_path:
            try:
                slide.shapes.add_movie(
                    speech_path,
                    Inches(8.5), Inches(0.5),
                    width=Inches(1.0), height=Inches(1.0),
                    mime_type='audio/mpeg'
                )
            except Exception:
                # Some PowerPoint versions don’t support add_movie()
                pass

    pptx_filename = "presentation.pptx"
    pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], pptx_filename)
    prs.save(pptx_path)
    print(f"✅ Saved presentation to {pptx_path}")
    return url_for('static', filename=f'outputs/{pptx_filename}')


# --- Gemini Functions ---

def generate_content_from_topic(topic):
    """
    Generate slides for a given topic.
    """
    print(f"Generating content for topic: {topic}")
    model = genai.GenerativeModel('models/gemini-2.5-flash')

    prompt = f"""
    You are an expert presentation creator. Create a 5-slide presentation about "{topic}".
    Return a valid JSON with structure:
    {{
      "slides": [
        {{ "title": "Slide Title", "content": "2–3 sentence narration" }}
      ]
    }}
    Ensure the JSON is syntactically correct.
    """
    try:
        response = model.generate_content(prompt)
        return clean_json_response(response.text)
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
        return None


def generate_content_from_text(text_content):
    """
    Summarizes provided text into slides.
    """
    print("Generating content from text...")
    model = genai.GenerativeModel('models/gemini-2.5-flash')

    prompt = f"""
    Convert the following text into a 5-slide presentation.
    Each slide must include:
    1. A catchy "title".
    2. A concise "content" of 2–3 sentences.

    Return a single valid JSON object like:
    {{
      "slides": [
        {{ "title": "Slide 1 Title", "content": "..." }},
        {{ "title": "Slide 2 Title", "content": "..." }},
        ...
      ]
    }}

    TEXT:
    ---
    {text_content}
    ---
    """
    try:
        response = model.generate_content(prompt)
        return clean_json_response(response.text)
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
        return None


def scrape_url(url):
    """
    Extracts readable text from a webpage.
    """
    print(f"Scraping: {url}")
    try:
        r = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
        if r.status_code != 200:
            return None, f"Failed to fetch URL (status {r.status_code})"

        soup = BeautifulSoup(r.text, 'html.parser')
        for tag in soup(['script', 'style']):
            tag.decompose()
        text = ' '.join(soup.get_text().split())
        return text, None if text else "No content found."
    except Exception as e:
        print(f"Error scraping URL: {e}")
        return None, str(e)


# --- Flask Routes ---

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
def handle_generate():
    try:
        topic = request.json.get('topic')
        if not topic:
            return jsonify({'error': 'No topic provided'}), 400

        data = generate_content_from_topic(topic)
        if not data or 'slides' not in data:
            return jsonify({'error': 'Failed to generate presentation.'}), 500

        download_url = create_powerpoint(data['slides'])
        return jsonify({'download_url': download_url})
    except Exception as e:
        print("Error in /generate:", e)
        return jsonify({'error': str(e)}), 500


@app.route('/generate-from-text', methods=['POST'])
def handle_generate_from_text():
    try:
        text = request.json.get('text')
        if not text:
            return jsonify({'error': 'No text provided'}), 400

        data = generate_content_from_text(text)
        if not data or 'slides' not in data:
            return jsonify({'error': 'Failed to summarize text.'}), 500

        download_url = create_powerpoint(data['slides'])
        return jsonify({'download_url': download_url})
    except Exception as e:
        print("Error in /generate-from-text:", e)
        return jsonify({'error': str(e)}), 500


@app.route('/generate-from-url', methods=['POST'])
def handle_generate_from_url():
    try:
        url = request.json.get('url')
        if not url:
            return jsonify({'error': 'No URL provided'}), 400

        text, err = scrape_url(url)
        if err:
            return jsonify({'error': err}), 500

        data = generate_content_from_text(text)
        if not data or 'slides' not in data:
            return jsonify({'error': 'Failed to summarize URL.'}), 500

        download_url = create_powerpoint(data['slides'])
        return jsonify({'download_url': download_url})
    except Exception as e:
        print("Error in /generate-from-url:", e)
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    app.run(debug=True)
